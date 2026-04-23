"""PPTX source ingestion: native text extraction + optional image rendering."""

from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

from core.intake.pdf_input import PdfConversionError, prepare_assets_from_pdf


class PptxConversionError(RuntimeError):
    """Raised when a PPTX source cannot be converted."""



def _extract_native_text(pptx_path: Path) -> list[dict[str, str]]:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise PptxConversionError(
            "No se pudo extraer texto nativo de PPTX porque falta dependencia `python-pptx`."
        ) from exc

    try:
        presentation = Presentation(str(pptx_path))
    except Exception as exc:
        raise PptxConversionError(f"No se pudo abrir el PPTX: {pptx_path.name}.") from exc

    slides_text: list[dict[str, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = str(shape.text).strip()
                if txt:
                    texts.append(txt)
        slides_text.append({"slide": index, "text": "\n".join(texts).strip()})
    return slides_text


def _convert_pptx_to_pdf(*, pptx_path: Path, temp_dir: Path) -> Path | None:
    soffice = shutil.which("soffice")
    libreoffice = shutil.which("libreoffice")
    office_bin = soffice or libreoffice

    if office_bin is None:
        return None

    temp_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        office_bin,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(temp_dir),
        str(pptx_path),
    ]
    process = subprocess.run(cmd, capture_output=True, text=True)
    if process.returncode != 0:
        details = (process.stderr or process.stdout or "").strip()
        raise PptxConversionError(
            "No se pudo convertir PPTX a PDF con LibreOffice. "
            f"Detalle: {details or 'sin detalle adicional.'}"
        )

    pdf_path = temp_dir / f"{pptx_path.stem}.pdf"
    return pdf_path if pdf_path.exists() else None


def prepare_assets_from_pptx(*, pptx_path: Path, destination_dir: Path, temp_dir: Path) -> dict:
    native_text_slides = _extract_native_text(pptx_path)

    warnings: list[str] = []
    prepared_images: list = []

    pdf_path = _convert_pptx_to_pdf(pptx_path=pptx_path, temp_dir=temp_dir)
    if pdf_path is None:
        warnings.append(
            "LibreOffice no disponible; se omite render de slides a imágenes. "
            "Se continúa con texto nativo de PPTX."
        )
    else:
        try:
            converted = prepare_assets_from_pdf(pdf_path=pdf_path, destination_dir=destination_dir)
            prepared_images = converted["prepared_images"]
        except PdfConversionError as exc:
            warnings.append(f"No se pudieron renderizar imágenes desde PPTX convertido: {exc}")

    return {
        "prepared_images": prepared_images,
        "native_text_entries": [
            {
                "index": item["slide"],
                "source": f"{pptx_path}#slide={item['slide']}",
                "text": item["text"],
                "kind": "pptx_slide",
            }
            for item in native_text_slides
        ],
        "warnings": warnings,
    }
